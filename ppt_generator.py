"""
PPT Generator Module - Creates PowerPoint presentations with text overlays.
Uses background images with programmatically added text for clear, editable content.
"""
import os
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from config import SLIDE_WIDTH, SLIDE_HEIGHT, OUTPUT_DIR


def rgb_color(r, g, b):
    """Create RGB color from r, g, b values (0-255)."""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


class PPTGenerator:
    """Generates PowerPoint presentations with background images and text overlays."""
    
    def __init__(self):
        self.width = Inches(SLIDE_WIDTH)
        self.height = Inches(SLIDE_HEIGHT)
        os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    def create_ppt_from_images(
        self, 
        image_paths: List[str], 
        output_filename: str,
        slides_data: Optional[List[Dict[str, Any]]] = None
    ) -> str:
        """
        Create a PPT with background images and text overlays.
        
        Args:
            image_paths: List of paths to background images
            output_filename: Output filename (without extension)
            slides_data: List of slide data with title and content
            
        Returns:
            Path to the created PPT file
        """
        prs = Presentation()
        prs.slide_width = self.width
        prs.slide_height = self.height
        
        blank_layout = prs.slide_layouts[6]  # Blank layout
        
        for i, image_path in enumerate(image_paths):
            if not os.path.exists(image_path):
                print(f"Warning: Image not found: {image_path}")
                continue
            
            # Create slide
            slide = prs.slides.add_slide(blank_layout)
            
            # Add background image (fill entire slide)
            slide.shapes.add_picture(
                image_path,
                left=0,
                top=0,
                width=self.width,
                height=self.height
            )
            
            # Get slide data
            slide_info = slides_data[i] if slides_data and i < len(slides_data) else {}
            title = slide_info.get('title', '')
            content = slide_info.get('content', '')
            slide_num = slide_info.get('slide_number', i + 1)
            
            # Add text overlays
            if slide_num == 1:
                # Cover slide - centered title
                self._add_cover_text(slide, title, content)
            elif slides_data and slide_num == len(slides_data):
                # Last slide - thank you page
                self._add_ending_text(slide, title)
            else:
                # Content slides
                self._add_content_text(slide, title, content)
        
        # Save the presentation
        output_path = os.path.join(OUTPUT_DIR, f"{output_filename}.pptx")
        prs.save(output_path)
        print(f"PPT saved to: {output_path}")
        
        return output_path
    
    def _add_cover_text(self, slide, title: str, subtitle: str = ""):
        """Add centered title and subtitle for cover slide."""
        # Semi-transparent overlay for better text readability
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=0,
            top=Inches(2.5),
            width=self.width,
            height=Inches(3)
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = rgb_color(0, 0, 0)
        # Set transparency
        overlay.fill.fore_color.brightness = 0.3
        overlay.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(2.8),
            width=Inches(SLIDE_WIDTH - 1),
            height=Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = rgb_color(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                left=Inches(0.5),
                top=Inches(4.2),
                width=Inches(SLIDE_WIDTH - 1),
                height=Inches(0.8)
            )
            sub_frame = sub_box.text_frame
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = subtitle.replace('；', ' · ').replace(';', ' · ')
            sub_para.font.size = Pt(24)
            sub_para.font.color.rgb = rgb_color(200, 200, 200)
            sub_para.alignment = PP_ALIGN.CENTER
    
    def _add_content_text(self, slide, title: str, content: str):
        """Add title and bullet points for content slides."""
        # Title background bar
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=0,
            top=0,
            width=self.width,
            height=Inches(1.4)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = rgb_color(0, 0, 0)
        title_bg.fill.fore_color.brightness = 0.2
        title_bg.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            left=Inches(0.6),
            top=Inches(0.35),
            width=Inches(SLIDE_WIDTH - 1.2),
            height=Inches(0.9)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = rgb_color(255, 255, 255)
        
        # Content area with semi-transparent background
        if content:
            content_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left=Inches(0.5),
                top=Inches(1.8),
                width=Inches(SLIDE_WIDTH - 1),
                height=Inches(4.8)
            )
            content_bg.fill.solid()
            content_bg.fill.fore_color.rgb = rgb_color(0, 0, 0)
            content_bg.fill.fore_color.brightness = 0.4
            content_bg.line.fill.background()
            
            # Content text
            content_box = slide.shapes.add_textbox(
                left=Inches(0.8),
                top=Inches(2.0),
                width=Inches(SLIDE_WIDTH - 1.6),
                height=Inches(4.4)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            # Parse content points (split by ; or ；)
            points = [p.strip() for p in content.replace('；', ';').split(';') if p.strip()]
            
            for idx, point in enumerate(points):
                if idx == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()
                
                para.text = f"• {point}"
                para.font.size = Pt(28)
                para.font.color.rgb = rgb_color(255, 255, 255)
                para.space_before = Pt(12)
                para.space_after = Pt(8)
    
    def _add_ending_text(self, slide, title: str):
        """Add centered text for ending/thank you slide."""
        # Overlay
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=0,
            top=Inches(2.8),
            width=self.width,
            height=Inches(2.5)
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = rgb_color(0, 0, 0)
        overlay.fill.fore_color.brightness = 0.3
        overlay.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(3.2),
            width=Inches(SLIDE_WIDTH - 1),
            height=Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = rgb_color(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER


def test_ppt_generator():
    """Test the PPT generator with sample slides."""
    import tempfile
    from PIL import Image
    
    # Create test background images
    test_images = []
    colors = [(52, 73, 94), (41, 128, 185), (39, 174, 96), (142, 68, 173), (231, 76, 60)]
    
    for i, color in enumerate(colors):
        img = Image.new('RGB', (1920, 1080), color)
        path = os.path.join(tempfile.gettempdir(), f"test_bg_{i+1}.png")
        img.save(path)
        test_images.append(path)
    
    # Sample slides data
    slides_data = [
        {
            "slide_number": 1,
            "title": "人工智能的未来",
            "content": "探索AI技术的无限可能"
        },
        {
            "slide_number": 2,
            "title": "核心技术",
            "content": "机器学习；深度学习；自然语言处理；计算机视觉"
        },
        {
            "slide_number": 3,
            "title": "应用场景",
            "content": "智能客服；自动驾驶；医疗诊断；金融风控"
        },
        {
            "slide_number": 4,
            "title": "发展趋势",
            "content": "多模态融合；边缘计算；可解释AI；通用人工智能"
        },
        {
            "slide_number": 5,
            "title": "感谢聆听",
            "content": ""
        }
    ]
    
    generator = PPTGenerator()
    output_path = generator.create_ppt_from_images(test_images, "test_overlay_ppt", slides_data)
    print(f"Test PPT created: {output_path}")
    
    # Cleanup
    for path in test_images:
        os.remove(path)


if __name__ == "__main__":
    test_ppt_generator()
